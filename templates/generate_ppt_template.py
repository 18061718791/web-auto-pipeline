#!/usr/bin/env python3
"""
generate_ppt_template.py — 程序化 PPT 生成模板

本模板提供配色常量 + 工具函数 + 13页标准结构，可直接填充内容生成 .pptx。

使用方法：
  1. 复制此文件为新脚本
  2. 修改标题/内容/截图路径
  3. python your_script.py → output.pptx

依赖: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ════════════════════ 配色常量 ════════════════════
BG_DARK      = RGBColor(0x0B, 0x19, 0x26)
CARD_BG      = RGBColor(0x13, 0x2B, 0x45)
ACCENT_GOLD  = RGBColor(0xF5, 0x9E, 0x0B)
ACCENT_CYAN  = RGBColor(0x08, 0x91, 0xB2)
TEXT_WHITE   = RGBColor(0xF0, 0xF6, 0xFC)
TEXT_GRAY    = RGBColor(0x8B, 0x94, 0x9E)
TEXT_GREEN   = RGBColor(0x3F, 0xB9, 0x50)
TEXT_RED     = RGBColor(0xF8, 0x51, 0x49)
BORDER_GRAY  = RGBColor(0x30, 0x36, 0x3D)
DIVIDER_CLR  = RGBColor(0x21, 0x2D, 0x3D)

# ════════════════════ 工具函数 ════════════════════

def set_bg(slide, color=BG_DARK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    """添加矩形并填充颜色"""
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, left, top, width, height, text, font_size=18,
             color=TEXT_WHITE, bold=False, align=PP_ALIGN.LEFT):
    """添加单段文本框"""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Microsoft YaHei'
    p.alignment = align
    return box

def add_bullets(slide, left, top, width, height, items, font_size=14,
                color=TEXT_GRAY, spacing=Pt(8)):
    """添加带圆点的列表项"""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Microsoft YaHei'
        p.space_after = spacing
    return box

def add_card(slide, left, top, width, height, title, icon="",
             fill=CARD_BG, border=BORDER_GRAY, title_color=ACCENT_GOLD):
    """添加带标题条的卡片"""
    card = add_rect(slide, left, top, width, height, fill, border)
    # 左侧竖条
    add_rect(slide, left, top, Inches(0.06), height, title_color)
    # 标题
    add_text(slide, left + Inches(0.2), top + Inches(0.08),
             width - Inches(0.3), Inches(0.4),
             f"{icon} {title}", font_size=16, color=title_color, bold=True)
    return card

def add_image_safe(slide, path, left, top, width, height=None):
    """安全添加图片（不存在时跳过）"""
    import os
    if not os.path.exists(path):
        return False
    kwargs = {"left": left, "top": top, "width": width}
    if height:
        kwargs["height"] = height
    slide.shapes.add_picture(path, **kwargs)
    return True

# ════════════════════ 13 页骨架 ════════════════════

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    W, H = prs.slide_width, prs.slide_height

    # ── 第1页: 封面 ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_rect(slide, 0, 0, W, Inches(0.08), ACCENT_GOLD)
    add_text(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(1.2),
             "TITLE HERE", font_size=48, color=TEXT_WHITE, bold=True)
    add_text(slide, Inches(0.8), Inches(2.7), Inches(11), Inches(0.8),
             "Subtitle Here", font_size=28, color=ACCENT_GOLD, bold=True)
    add_text(slide, Inches(0.8), Inches(3.6), Inches(11), Inches(0.8),
             "Tagline / 一句话介绍", font_size=16, color=TEXT_GRAY)
    add_rect(slide, 0, Inches(7.2), W, Inches(0.3), CARD_BG)

    # ── 第2页: 痛点 ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_rect(slide, 0, 0, W, Inches(0.06), ACCENT_GOLD)
    add_text(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
             "痛点引入", font_size=30, color=TEXT_WHITE, bold=True)
    add_text(slide, Inches(0.6), Inches(0.85), Inches(10), Inches(0.4),
             "副标题说明", font_size=14, color=TEXT_GRAY)
    # 每痛点: 红色竖条 + 标题 + 描述
    pains = [("痛点标题", "痛点描述"), ("痛点标题", "痛点描述")]
    for i, (title, desc) in enumerate(pains):
        y = Inches(1.4) + Inches(1.4) * i
        add_rect(slide, Inches(0.6), y, Inches(0.06), Inches(1.2), TEXT_RED)
        add_text(slide, Inches(0.85), y, Inches(11.5), Inches(0.4),
                 title, font_size=16, color=TEXT_WHITE, bold=True)
        add_text(slide, Inches(0.85), y + Inches(0.4), Inches(11.5), Inches(0.7),
                 desc, font_size=12, color=TEXT_GRAY)

    # ── 第3页: 方案映射 ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_text(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
             "痛点 → 方案", font_size=30, color=TEXT_WHITE, bold=True)
    # 每行: 痛点卡片 | → | 方案卡片

    # ── 第4页: 演进脉络 ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_text(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
             "演进脉络", font_size=30, color=TEXT_WHITE, bold=True)
    # 时间线: 横线 + 圆点 + 卡片

    # ── 第5-9页: 核心能力（用户自定义） ──
    for pi in range(5):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(slide)
        add_text(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
                 f"核心能力 {pi+1}", font_size=30, color=TEXT_WHITE, bold=True)

    # ── 第10页: 优势对比 ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_text(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
             "核心优势对比", font_size=30, color=TEXT_WHITE, bold=True)
    # 用 add_rect 行代替 add_table 以获得更灵活的颜色控制

    # ── 第11页: 实践案例 ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_text(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
             "实践案例", font_size=30, color=TEXT_WHITE, bold=True)

    # ── 第12页: 核心价值 ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_text(slide, Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
             "核心价值", font_size=30, color=TEXT_WHITE, bold=True)

    # ── 第13页: Thank You ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_text(slide, Inches(0.8), Inches(2.0), Inches(11), Inches(1.0),
             "Thank You", font_size=48, color=TEXT_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.8), Inches(3.0), Inches(11), Inches(0.6),
             "Q & A", font_size=28, color=ACCENT_GOLD, bold=True, align=PP_ALIGN.CENTER)

    return prs


if __name__ == "__main__":
    prs = create_presentation()
    output = "output_presentation.pptx"
    prs.save(output)
    print(f"OK: {output} ({len(prs.slides)} pages)")
